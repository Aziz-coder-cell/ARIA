from pptx import Presentation
from pptx.util import Inches,Pt
from ai import ppt_details
import json

def generate_ppt(topic,num_slides):
    prs = Presentation()
    prompt = f'Give me bullet points for {topic}, number of slides = {num_slides}, each slide should contain 4 points only. Return ONLY a JSON array, no markdown, no extra text: [{{"title": "Slide title", "bullets": ["point 1", "point 2"]}}]'
    result = ppt_details(prompt)
    data = json.loads(result)

    for i in range(num_slides):
        slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(slide_layout)
        title = slide.shapes.title 
        title.text = data[i]["title"]
        body = slide.placeholders[1]
        tf = body.text_frame
        tf.text = data[i]["bullets"][0]
        for j in range(1,4):
            p = tf.add_paragraph()
            p.text = data[i]["bullets"][j]

    prs.save(f"{topic}.pptx")
    print("The file is saved in the same folder where all the files are present")